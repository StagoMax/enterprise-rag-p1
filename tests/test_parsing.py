from __future__ import annotations

from pathlib import Path

import pytest

from enterprise_rag import parsing
from enterprise_rag.models import DocumentStatus
from enterprise_rag.parsing import (
    BlockKind,
    DoclingDocumentParser,
    ParseBackend,
    parse_documents,
    suppress_header_footer,
    to_document_input,
)

# 默认解析器优先 Docling，而 Docling 的 PDF 流水线需要模型权重（可能触发下载）。
# 断言用的解析器固定走内置抽取器，保证无 GPU、无网络也可复现。
FALLBACK_PARSER = DoclingDocumentParser(prefer_docling=False)

TABLE_TOKEN = "GATE-7788"
HEADER_LINE = "Orion Confidential Internal Only"


# ---------------------------------------------------------------------------
# 夹具文件在测试期生成，不入库任何二进制
# ---------------------------------------------------------------------------


def _write_docx(path: Path) -> Path:
    docx = pytest.importorskip("docx")
    path.parent.mkdir(parents=True, exist_ok=True)
    document = docx.Document()
    document.add_heading("1. VPN 接入策略", level=1)
    document.add_paragraph("员工访问生产只读控制台前必须连接 corp-shanghai VPN。")
    document.add_heading("1.1 常见错误码", level=2)
    document.add_paragraph("错误码 VPN-401 表示设备证书尚未登记。", style="List Bullet")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "门禁编号"
    table.cell(0, 1).text = "阈值"
    table.cell(1, 0).text = TABLE_TOKEN
    table.cell(1, 1).text = "0.92"
    document.save(str(path))
    return path


def _write_pptx(path: Path) -> Path:
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "发布门禁规范"
    slide.placeholders[1].text = "路由回归必须通过\n权限测试必须零越权"
    second = presentation.slides.add_slide(presentation.slide_layouts[5])
    second.shapes.title.text = "门禁指标"
    table = second.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(6), Inches(1)).table
    table.cell(0, 0).text = "门禁编号"
    table.cell(0, 1).text = "阈值"
    table.cell(1, 0).text = TABLE_TOKEN
    table.cell(1, 1).text = "0.92"
    presentation.save(str(path))
    return path


def _write_xlsx(path: Path) -> Path:
    openpyxl = pytest.importorskip("openpyxl")
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "门禁指标"
    sheet.append(["门禁编号", "阈值"])
    sheet.append([TABLE_TOKEN, "0.92"])
    workbook.save(str(path))
    return path


def _write_html(path: Path) -> Path:
    path.write_text(
        """<!doctype html>
<html><head><title>Orion 发布门禁规范</title><style>p{color:red}</style></head>
<body>
<header>Orion Confidential Internal Only</header>
<nav>首页 / 规范 / 门禁</nav>
<h1>1. 发布门禁</h1>
<p>发布前必须通过路由、引用、权限与 SQL 安全回归测试。</p>
<h2>1.1 指标</h2>
<ul><li>路由正确率不低于 90%</li></ul>
<table>
<tr><th>门禁编号</th><th>阈值</th></tr>
<tr><td>GATE-7788</td><td>0.92</td></tr>
</table>
<script>console.log("noise")</script>
<footer>第 1 页 共 1 页</footer>
</body></html>
""",
        encoding="utf-8",
    )
    return path


def _write_markdown(path: Path) -> Path:
    path.write_text(
        "# 发布门禁规范\n\n"
        "发布前必须通过四类回归测试。\n\n"
        "## 指标\n\n"
        "- 路由正确率不低于 90%\n"
        "- 权限测试必须零越权\n",
        encoding="utf-8",
    )
    return path


def _write_pdf(path: Path, pages: list[list[str]]) -> Path:
    fpdf = pytest.importorskip("fpdf")
    from fpdf.enums import XPos, YPos

    pdf = fpdf.FPDF()
    pdf.set_auto_page_break(auto=False)
    pdf.set_font("Helvetica", size=11)
    for lines in pages:
        pdf.add_page()
        for line in lines:
            pdf.cell(0, 7, line, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    pdf.output(str(path))
    return path


def _write_blank_pdf(path: Path, pages: int = 1) -> Path:
    fpdf = pytest.importorskip("fpdf")
    pdf = fpdf.FPDF()
    for _ in range(pages):
        pdf.add_page()
    pdf.output(str(path))
    return path


def _table_rows(parsed: parsing.ParsedDocument) -> list[list[str]]:
    return [row for table in parsed.tables for row in table.rows]


# ---------------------------------------------------------------------------
# 逐格式抽取
# ---------------------------------------------------------------------------


def test_docx_extracts_structure_and_separates_tables(tmp_path: Path) -> None:
    parsed = FALLBACK_PARSER.parse(_write_docx(tmp_path / "vpn.docx"))

    assert not parsed.failed
    assert parsed.backend is ParseBackend.FALLBACK
    assert "corp-shanghai VPN" in parsed.text
    assert "VPN-401" in parsed.text
    kinds = {block.kind for block in parsed.blocks}
    assert BlockKind.HEADING in kinds
    assert BlockKind.LIST in kinds
    # 表格必须走关系库路线，不能被摊平进散文
    assert TABLE_TOKEN not in parsed.text
    assert len(parsed.tables) == 1
    assert [TABLE_TOKEN, "0.92"] in parsed.tables[0].rows
    assert parsed.tables[0].header == ["门禁编号", "阈值"]


def test_docx_sections_follow_document_numbering(tmp_path: Path) -> None:
    parsed = FALLBACK_PARSER.parse(_write_docx(tmp_path / "vpn.docx"))
    sections = [block.section for block in parsed.blocks if block.kind is BlockKind.HEADING]
    assert sections == ["1", "1.1"]
    assert any(block.heading_path == ["1. VPN 接入策略"] for block in parsed.blocks)


def test_pptx_extracts_per_slide_pages_and_tables(tmp_path: Path) -> None:
    parsed = FALLBACK_PARSER.parse(_write_pptx(tmp_path / "gate.pptx"))

    assert not parsed.failed
    assert parsed.page_count == 2
    assert "路由回归必须通过" in parsed.text
    assert {block.page for block in parsed.blocks} == {1, 2}
    assert TABLE_TOKEN not in parsed.text
    assert [TABLE_TOKEN, "0.92"] in _table_rows(parsed)
    assert parsed.tables[0].page == 2


def test_xlsx_sheets_become_tables_not_prose(tmp_path: Path) -> None:
    parsed = FALLBACK_PARSER.parse(_write_xlsx(tmp_path / "metrics.xlsx"))

    assert not parsed.failed
    assert "门禁指标" in parsed.text
    assert TABLE_TOKEN not in parsed.text
    assert parsed.tables[0].caption == "门禁指标"
    assert parsed.tables[0].rows == [["门禁编号", "阈值"], [TABLE_TOKEN, "0.92"]]


def test_html_drops_boilerplate_scripts_and_keeps_table(tmp_path: Path) -> None:
    parsed = FALLBACK_PARSER.parse(_write_html(tmp_path / "gate.html"))

    assert not parsed.failed
    assert parsed.title == "Orion 发布门禁规范"
    assert "SQL 安全回归测试" in parsed.text
    assert "console.log" not in parsed.text
    assert HEADER_LINE not in parsed.text
    assert "第 1 页 共 1 页" not in parsed.text
    assert any("首页" in dropped for dropped in parsed.suppressed_lines)
    assert TABLE_TOKEN not in parsed.text
    assert [TABLE_TOKEN, "0.92"] in _table_rows(parsed)


def test_markdown_headings_and_lists(tmp_path: Path) -> None:
    parsed = FALLBACK_PARSER.parse(_write_markdown(tmp_path / "gate.md"))

    assert not parsed.failed
    assert parsed.title == "发布门禁规范"
    headings = [block.text for block in parsed.blocks if block.kind is BlockKind.HEADING]
    assert headings == ["发布门禁规范", "指标"]
    assert sum(block.kind is BlockKind.LIST for block in parsed.blocks) == 2


def test_pdf_text_extraction_suppresses_header_and_page_numbers(tmp_path: Path) -> None:
    bodies = [
        "Release gate policy applies to every Orion P1 deployment.",
        "Routing, citation, permission and SQL safety regressions must pass.",
        "Any unauthorized chunk recall blocks the release immediately.",
    ]
    pages = [
        [HEADER_LINE, body, f"Page {number} of {len(bodies)}"]
        for number, body in enumerate(bodies, start=1)
    ]
    parsed = FALLBACK_PARSER.parse(_write_pdf(tmp_path / "policy.pdf", pages))

    assert not parsed.failed, parsed.error
    assert parsed.warnings == [], parsed.warnings
    assert parsed.page_count == 3
    assert parsed.needs_ocr is False
    assert "Release gate policy" in parsed.text
    assert "SQL safety regressions" in parsed.text
    assert HEADER_LINE not in parsed.text
    assert "Page 2 of 3" not in parsed.text
    assert HEADER_LINE in parsed.suppressed_lines
    assert {block.page for block in parsed.blocks} == {1, 2, 3}


# ---------------------------------------------------------------------------
# 锚点稳定性
# ---------------------------------------------------------------------------


def test_anchors_are_stable_and_path_independent(tmp_path: Path) -> None:
    first_path = _write_docx(tmp_path / "run-a" / "vpn.docx")
    second_path = tmp_path / "run-b" / "renamed.docx"
    second_path.parent.mkdir(parents=True, exist_ok=True)
    second_path.write_bytes(first_path.read_bytes())

    first = DoclingDocumentParser(prefer_docling=False).parse(first_path)
    second = DoclingDocumentParser(prefer_docling=False).parse(second_path)
    third = FALLBACK_PARSER.parse(first_path)

    assert first.anchors() == second.anchors() == third.anchors()
    assert first.anchors(), "解析结果必须带锚点，引用链依赖它"
    assert len(set(first.anchors())) == len(first.anchors())


def test_html_anchors_stable_across_runs(tmp_path: Path) -> None:
    path = _write_html(tmp_path / "gate.html")
    assert FALLBACK_PARSER.parse(path).anchors() == FALLBACK_PARSER.parse(path).anchors()


# ---------------------------------------------------------------------------
# 扫描件与 OCR 路由
# ---------------------------------------------------------------------------


def test_text_free_pdf_is_flagged_as_needing_ocr(tmp_path: Path) -> None:
    path = _write_blank_pdf(tmp_path / "scan.pdf", pages=2)
    parser = DoclingDocumentParser(prefer_docling=False, enable_ocr=False)
    parsed = parser.parse(path)

    assert parsed.needs_ocr is True
    assert parsed.ocr_applied is False
    # 扫描件不是解析失败，只是需要 OCR 兜底
    assert parsed.failed is False
    assert parsed.has_prose is False
    assert any("OCR" in warning for warning in parsed.warnings)


def test_scanned_pdf_routes_to_ocr(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    path = _write_blank_pdf(tmp_path / "scan.pdf", pages=2)
    monkeypatch.setattr(
        parsing,
        "ocr_pdf_pages",
        lambda target, **kwargs: (["1. 扫描件标题", "OCR 还原的第二页正文内容。"], []),
    )
    parsed = DoclingDocumentParser(prefer_docling=False).parse(path)

    assert parsed.needs_ocr is True
    assert parsed.ocr_applied is True
    assert parsed.backend is ParseBackend.OCR
    assert "OCR 还原的第二页正文内容。" in parsed.text
    assert parsed.degraded is True


def test_ocr_dependency_gap_is_reported_not_raised(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    # OCR 引擎缺失时必须返回原因而不是抛异常（不联网、不下载模型）
    monkeypatch.setattr(parsing, "_load_ocr_reader", lambda: (None, "paddleocr 不可用"))
    pages, warnings = parsing.ocr_pdf_pages(_write_blank_pdf(tmp_path / "scan.pdf"))

    assert pages == []
    assert warnings and "OCR 跳过" in warnings[0]


class _RapidOcrV3Result:
    """RapidOCR 3.x 返回带 txts 属性的结果对象。"""

    txts = ("门禁编号 GATE-7788", "阈值 0.92")


@pytest.mark.parametrize(
    ("raw", "expected"),
    [
        (_RapidOcrV3Result(), ["门禁编号 GATE-7788", "阈值 0.92"]),
        ([{"rec_texts": ["第一行", "第二行"]}], ["第一行", "第二行"]),
        (([[[0, 0], [1, 1]], "第一行", 0.98],), ["第一行"]),
        ("unexpected payload", []),
        (([[[[0, 0]], "第一行", 0.98], [[[1, 1]], "第二行", 0.9]], 0.5), ["第一行", "第二行"]),
        ([[[[[0, 0]], ("第一行", 0.98)], [[[1, 1]], ("第二行", 0.9)]]], ["第一行", "第二行"]),
        (None, []),
    ],
)
def test_recognized_lines_normalizes_engine_shapes(raw: object, expected: list[str]) -> None:
    assert parsing._recognized_lines(raw) == expected


def test_ocr_disabled_pdf_keeps_needs_ocr_without_touching_engines(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    def _fail() -> tuple[None, str]:
        raise AssertionError("enable_ocr=False 时不应加载 OCR 引擎")

    monkeypatch.setattr(parsing, "_load_ocr_reader", _fail)
    parser = DoclingDocumentParser(prefer_docling=False, enable_ocr=False)
    parsed = parser.parse(_write_blank_pdf(tmp_path / "scan.pdf"))

    assert parsed.needs_ocr is True
    assert parsed.ocr_applied is False


# ---------------------------------------------------------------------------
# 优雅降级
# ---------------------------------------------------------------------------


def test_corrupt_docx_degrades_instead_of_raising(tmp_path: Path) -> None:
    path = tmp_path / "broken.docx"
    path.write_bytes(b"this is definitely not an OOXML package")

    parsed = FALLBACK_PARSER.parse(path)

    assert parsed.failed is True
    assert parsed.backend is ParseBackend.NONE
    assert parsed.error is not None
    assert parsed.text == ""
    assert parsed.degraded is True


def test_corrupt_file_degrades_with_docling_preferred(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    monkeypatch.setattr(parsing, "_docling_import_error", lambda: "docling 不可用(测试注入)")
    path = tmp_path / "broken.docx"
    path.write_bytes(b"not ooxml at all")

    parsed = parsing.parse_document(path)

    assert parsed.failed is True
    assert parsed.error is not None
    assert any("docling" in warning for warning in parsed.warnings)


def test_unsupported_suffix_and_missing_file_are_data_not_exceptions(tmp_path: Path) -> None:
    unsupported = tmp_path / "archive.rar"
    unsupported.write_bytes(b"\x00\x01")
    unsupported_result = FALLBACK_PARSER.parse(unsupported)
    assert unsupported_result.failed is True
    assert "不受支持" in (unsupported_result.error or "")
    assert FALLBACK_PARSER.supports(unsupported) is False

    missing = FALLBACK_PARSER.parse(tmp_path / "nope.pdf")
    assert missing.failed is True
    assert "不存在" in (missing.error or "")


def test_parse_documents_surfaces_per_file_failures(tmp_path: Path) -> None:
    good = _write_markdown(tmp_path / "gate.md")
    broken = tmp_path / "broken.xlsx"
    broken.write_bytes(b"garbage")

    results = parse_documents([good, broken], parser=FALLBACK_PARSER)

    assert [result.failed for result in results] == [False, True]
    assert results[1].error is not None


def test_suppress_header_footer_keeps_body_lines() -> None:
    pages = [f"机密 内部资料\n正文第 {index} 段内容。\n{index}" for index in range(1, 4)]
    cleaned, suppressed = suppress_header_footer(pages)

    assert all("正文第" in page for page in cleaned)
    assert "机密 内部资料" in suppressed
    assert all("机密 内部资料" not in page for page in cleaned)
    assert "1" in suppressed


# ---------------------------------------------------------------------------
# 接入既有管线
# ---------------------------------------------------------------------------


def test_to_document_input_respects_metadata_contract(tmp_path: Path) -> None:
    parsed = FALLBACK_PARSER.parse(_write_docx(tmp_path / "vpn.docx"))

    document = to_document_input(
        parsed,
        document_id="vpn-access-guide",
        owner="it-platform",
        business_class="technical-guide",
        allowed_roles={"engineering", " operations "},
        sensitivity="internal",
        version="1.2",
        source_uri="demo://it/vpn-access-guide",
    )

    assert document.document_id == "vpn-access-guide"
    assert document.title == "1. VPN 接入策略"
    assert document.owner == "it-platform"
    assert document.business_class == "technical-guide"
    assert document.allowed_roles == {"engineering", "operations"}
    assert document.version == "1.2"
    assert document.status is DocumentStatus.ACTIVE
    assert document.sensitivity == "internal"
    assert document.source_uri == "demo://it/vpn-access-guide"
    # 默认不摊平表格，只留可引用的锚点占位行
    assert TABLE_TOKEN not in document.content
    assert parsed.tables[0].anchor in document.content


def test_to_document_input_flows_into_existing_chunking(tmp_path: Path) -> None:
    from enterprise_rag.chunking import build_document, chunk_document

    parsed = FALLBACK_PARSER.parse(_write_markdown(tmp_path / "gate.md"))
    document = to_document_input(
        parsed,
        document_id="release-gate-policy",
        owner="engineering-enablement",
        business_class="engineering-policy",
        allowed_roles={"engineering"},
    )
    chunks = chunk_document(build_document(document))

    assert chunks
    assert chunks[0].document_id == "release-gate-policy"
    assert chunks[0].allowed_roles == frozenset({"engineering"})


def test_to_document_input_can_inline_tables_when_requested(tmp_path: Path) -> None:
    parsed = FALLBACK_PARSER.parse(_write_xlsx(tmp_path / "metrics.xlsx"))

    document = to_document_input(
        parsed,
        document_id="gate-metrics",
        owner="engineering-enablement",
        business_class="engineering-policy",
        allowed_roles={"engineering"},
        include_tables=True,
    )

    assert TABLE_TOKEN in document.content
    assert "| 门禁编号 | 阈值 |" in document.content
    assert document.source_uri is not None
    assert document.source_uri.startswith("file:")


def test_to_document_input_rejects_failed_parse(tmp_path: Path) -> None:
    path = tmp_path / "broken.pptx"
    path.write_bytes(b"garbage")
    parsed = FALLBACK_PARSER.parse(path)

    with pytest.raises(ValueError, match="解析失败"):
        to_document_input(
            parsed,
            document_id="broken",
            owner="owner",
            business_class="engineering-policy",
            allowed_roles={"engineering"},
        )


# ---------------------------------------------------------------------------
# Docling 适配层
#
# 这里用伪 Docling 文档驱动适配器，而不是启动真正的 DocumentConverter：真实转换器
# 会加载 torch/CUDA 与识别模型（PDF 流水线还需从 HuggingFace 下载权重），既联网又
# 会和进程内其他原生扩展抢资源。真实 Docling 后端另行手工验证。
# ---------------------------------------------------------------------------


class _FakeProvenance:
    def __init__(self, page_no: int) -> None:
        self.page_no = page_no


class _FakeItem:
    def __init__(self, label: str, text: str, page: int | None = None) -> None:
        self.label = label
        self.text = text
        self.prov = [_FakeProvenance(page)] if page else []


class _FakeCell:
    def __init__(self, text: str) -> None:
        self.text = text


class _FakeTableData:
    def __init__(self, rows: list[list[str]]) -> None:
        self.grid = [[_FakeCell(cell) for cell in row] for row in rows]


class _FakeTableItem:
    label = "table"
    text = ""

    def __init__(self, rows: list[list[str]], page: int) -> None:
        self.data = _FakeTableData(rows)
        self.prov = [_FakeProvenance(page)]

    def caption_text(self, document: object) -> str:
        # docling-core 里 caption_text 是方法而非属性
        return "门禁指标"


class _FakeDoclingDocument:
    def __init__(self) -> None:
        self.tables = [_FakeTableItem([["门禁编号", "阈值"], [TABLE_TOKEN, "0.92"]], 2)]
        self._items: list[tuple[object, int]] = [
            (_FakeItem("title", "1. 发布门禁", 1), 1),
            (_FakeItem("page_header", HEADER_LINE, 1), 1),
            (_FakeItem("text", "发布前必须通过四类回归测试。", 1), 1),
            (_FakeItem("list_item", "路由正确率不低于 90%", 1), 1),
            (_FakeItem("page_footer", "第 1 页 共 2 页", 1), 1),
            (self.tables[0], 1),
        ]

    def iterate_items(self) -> object:
        return iter(self._items)


class _FakeConverter:
    def __init__(self, document: object) -> None:
        self._document = document

    def convert(self, source: str) -> object:
        return type("_Result", (), {"document": self._document})()


def _use_fake_docling(monkeypatch: pytest.MonkeyPatch, document: object) -> None:
    monkeypatch.setattr(parsing, "_docling_import_error", lambda: "")
    monkeypatch.setattr(
        DoclingDocumentParser,
        "_converter_instance",
        lambda self: _FakeConverter(document),
    )


def test_docling_adapter_maps_blocks_pages_and_tables(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    _use_fake_docling(monkeypatch, _FakeDoclingDocument())
    parsed = parsing.parse_document(_write_docx(tmp_path / "vpn.docx"))

    assert parsed.backend is ParseBackend.DOCLING
    assert parsed.warnings == []
    assert parsed.degraded is False
    assert "发布前必须通过四类回归测试。" in parsed.text
    # 页眉、页脚由 Docling 标签直接丢弃
    assert HEADER_LINE not in parsed.text
    assert "第 1 页 共 2 页" not in parsed.text
    headings = [block for block in parsed.blocks if block.kind is BlockKind.HEADING]
    assert [block.section for block in headings] == ["1"]
    assert any(block.kind is BlockKind.LIST for block in parsed.blocks)
    assert {block.page for block in parsed.blocks} == {1}
    # 表格仍旧独立于散文，并能取到方法形态的 caption
    assert TABLE_TOKEN not in parsed.text
    assert parsed.tables[0].caption == "门禁指标"
    assert parsed.tables[0].page == 2
    assert parsed.tables[0].rows == [["门禁编号", "阈值"], [TABLE_TOKEN, "0.92"]]


def test_docling_anchors_match_across_runs(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    path = _write_docx(tmp_path / "vpn.docx")
    _use_fake_docling(monkeypatch, _FakeDoclingDocument())
    first = DoclingDocumentParser().parse(path)
    _use_fake_docling(monkeypatch, _FakeDoclingDocument())
    second = DoclingDocumentParser().parse(path)

    assert first.anchors() == second.anchors()


def test_docling_failure_degrades_with_recorded_reason(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    def _explode(self: DoclingDocumentParser) -> object:
        raise RuntimeError("模型权重缺失")

    monkeypatch.setattr(parsing, "_docling_import_error", lambda: "")
    monkeypatch.setattr(DoclingDocumentParser, "_converter_instance", _explode)

    parsed = DoclingDocumentParser().parse(_write_docx(tmp_path / "vpn.docx"))

    assert parsed.backend is ParseBackend.FALLBACK
    assert parsed.failed is False
    assert parsed.has_prose, "Docling 失败后仍必须能靠内置抽取器出正文"
    assert any("模型权重缺失" in warning for warning in parsed.warnings)


def test_docling_empty_output_degrades_to_fallback(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    class _EmptyDocument:
        tables: list[object] = []

        def iterate_items(self) -> object:
            return iter(())

    _use_fake_docling(monkeypatch, _EmptyDocument())
    parsed = DoclingDocumentParser().parse(_write_docx(tmp_path / "vpn.docx"))

    assert parsed.backend is ParseBackend.FALLBACK
    assert parsed.has_prose
    assert any("未产出任何块" in warning for warning in parsed.warnings)
